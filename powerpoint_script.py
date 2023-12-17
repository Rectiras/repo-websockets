from pptx import Presentation
from pptx.util import Inches

# This is the additional code I have used to generate my presentation quickly,
# then I have added design and adjusted it according to my personal preferences.
# If you would like to use this library you might need to:
# pip install python-pptx

def add_title_slide(prs, title):
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title

def add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Content Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = title
    content_box.text = content

def main():
    prs = Presentation()

    # Title Slide
    add_title_slide(prs, "WebSocket Implementation Comparison")

    # Latency Slide
    add_content_slide(
        prs,
        "Latency",
        "- Websockets: Average Latency - 0.000876 seconds\n"
        "- Tornado: Average Latency - 0.000965 seconds\n\n"
        "Observation: Websockets outperformed Tornado by approximately 10% in terms of latency."
    )

    # Throughput Slide
    add_content_slide(
        prs,
        "Throughput",
        "- Websockets: Throughput - 30.3015221337059 messages/second\n"
        "- Tornado: Throughput - 30.2804844366243 messages/second\n\n"
        "Observation: Throughput values between Websockets and Tornado are almost identical."
    )

    # Scalability Slide
    add_content_slide(
        prs,
        "Scalability",
        "Websockets:\n"
        "- High concurrency support\n"
        "- Asynchronous architecture for efficiency\n\n"
        "Tornado:\n"
        "- Well-suited for handling concurrent connections\n"
        "- May require fine-tuning for specific scalability requirements"
    )

    # Additional Metrics Slide
    add_content_slide(
        prs,
        "Additional Metrics",
        "Websockets:\n"
        "- Straightforward API\n"
        "- Active community and extensive documentation\n\n"
        "Tornado:\n"
        "- Rich feature set\n"
        "- Seamless integration into existing web applications"
    )

    # Conclusion Slide
    add_content_slide(
        prs,
        "Conclusion",
        "Considering the comparison in terms of latency and throughput:\n"
        "- Latency: Websockets demonstrated superior performance.\n"
        "- Throughput: Both libraries achieved nearly identical throughput, with a slight edge to Websockets.\n"
        "- Scalability: Websockets provides high concurrency support and asynchronous architecture.\n"
        "- Additional Metrics: Websockets offers a straightforward API and active community support.\n\n"
        "In conclusion, for this scenario, Websockets is recommended for its superior latency and comparable throughput."
    )

    # Save the presentation
    prs.save("WebSocket_Comparison_Presentation.pptx")

if __name__ == "__main__":
    main()
